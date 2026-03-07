
# Perform Cross-Source Deduplication:
#       Identify and merge near-duplicate content across different platforms.
#       Use techniques like MinHash to find duplicates even if they aren't exact matches.

"""
In a RAG (Retrieval-Augmented Generation) architecture, deduplication is typically applied during the Ingestion Pipeline (before embedding) or the Post-Retrieval phase (re-ranking).
Here are the top 10 implementation scenarios where this is critical:

   1. Multi-Channel Customer Support: When syncing data from Zendesk tickets, email chains, and Slack logs,
            the same solution is often repeated across platforms. Deduplication prevents the LLM
            from getting "looping" or redundant context.
   2. Enterprise Policy Hubs: Large companies often have Global, Regional, and Local versions of an HR policy.
            Deduplication ensures the RAG system prioritizes the most specific or latest version rather than
            retrieving three nearly identical documents.
   3. Medical Research & Clinical Trials: A single study may appear in PubMed, arXiv, and a private database.
            Deduplication prevents the model from "double-counting" evidence, which could lead to a false sense
            of scientific consensus.
   4. Software Documentation & Versioning: In a RAG system for developers, documentation for v1.0, v1.1, and v1.2
            will be 95% identical. Deduplication prevents the model from accidentally retrieving outdated syntax
            from an older version.
   5. Financial Audit & Compliance: Auditors use RAG to query thousands of invoices and contracts.
            Duplicate entries (e.g., a scanned copy and a digital PDF of the same bill) must be merged to
            ensure accurate financial totals.
   6. Legal Discovery (e-Discovery): During litigation, RAG systems scan through "Email Threads."
            Since each reply contains the entire previous history, deduplication (specifically "thread pruning")
            is required to avoid providing the same text 20 times in one context window.
   7. Product Catalog Q&A: In e-commerce, the same product might be listed in multiple categories with slight
            variations in descriptions. Deduplication ensures the chatbot provides a clean,
            singular answer about product specs.
   8. News Aggregation & Summarization: A RAG-based news bot might pull from Reuters, AP, and local outlets.
            MinHash is used here to identify that they are all reporting the same event,
            allowing the LLM to synthesize a single summary.
   9. Mergers & Acquisitions (M&A) Intelligence: When two companies merge, their Knowledge Bases often overlap
            (e.g., both have a "Travel Expense Policy"). Deduplication is used to "de-conflict" these sources
            during the data migration into a shared RAG.
   10. Academic Study Assistants: For students uploading lecture notes, slide decks, and textbook chapters,
            the same definitions often appear in all three. Deduplication optimizes the Context Window
            by removing these redundant explanations.
"""

from datasketch import MinHash, MinHashLSH
import re
import pandas as pd
from docx import Document
from pptx import Presentation

base_path = "..\\data\\stock_market\\"
files_to_process = ["Infosys_Q3_FY26_Full_Dataset.xlsx",
                   "Infosys_Q3_FY26_Management_Report.docx",
                   "Infosys_Q3_FY26_Investor_Presentation.pptx"]

# xlsx data collection
all_sheets = pd.read_excel(base_path + files_to_process[0], sheet_name=None)
my_xl_data_frame_combine = pd.concat(all_sheets.values(), ignore_index=True)

# pptx data collection
my_doc = Document(base_path + files_to_process[1])
my_doc_text = "\n".join([para.text for para in my_doc.paragraphs])

# docx data collection
my_ppt = Presentation(base_path + files_to_process[2])
my_ppt_text_lst = []
for each_slide in my_ppt.slides:
    for shape in each_slide.shapes:
        if hasattr(shape, "text"):  # Check if shape contains text
            my_ppt_text_lst.append(shape.text)
my_ppt_text =  "\n".join(my_ppt_text_lst)

# 1. Sample Data from different "Sources"
text_data = {"source_A": my_xl_data_frame_combine.to_string(),
             "source_B": my_doc_text,
             "source_C": my_ppt_text}


# # 1. Sample Data from different "Sources"
# text_data = {
#     "source_A": "The quick brown fox jumps over the lazy dog.",
#     "source_B": "The quick brown fox leaped over a lazy dog!",  # Near duplicate
#     "source_C": "A completely different sentence about cats.",  # Unique
#     "source_D": " quick brown foxes jumps over the lazy  dog"  # Near duplicate (missing period)
# }


def preprocess(text):
    """Simple tokenizer to create shingles."""
    text = text.lower()
    text = re.sub(r'[^\w\s]', '', text)
    return set(text.split())


# 2. Initialize LSH - Locality Sensitive Hashing
# threshold=0.8 means we want documents that are 80% similar
lsh = MinHashLSH(threshold=0.8, num_perm=128)

# 3. Create MinHashes and add to LSH index
min_hashes = {}
for doc_id, text in text_data.items():
    m = MinHash(num_perm=128)
    shingles = preprocess(text)
    for s in shingles:
        m.update(s.encode('utf8'))

    lsh.insert(doc_id, m)
    min_hashes[doc_id] = m

# 4. Query for duplicates
print("--- Deduplication Results ---")
processed = set()

for doc_id in text_data.keys():
    if doc_id in processed:
        continue

    # Find all items similar to this one
    result = lsh.query(min_hashes[doc_id])

    if len(result) > 1:
        print(f"Cluster found: {result}")
        # In a real merge, you'd pick one 'primary' and discard/link the others
        for r in result:
            processed.add(r)
    else:
        print(f"Unique item: {doc_id}")
